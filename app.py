import streamlit as st
import fitz  # PyMuPDF
import openai
import tempfile
import os
import docx
from pptx import Presentation
from docx import Document
from pptx import Presentation as PptPresentation
from fpdf import FPDF
import pandas as pd

# ---- CONFIG ----
openai.api_key = st.secrets["OPENAI_API_KEY"]
MODEL = "gpt-3.5-turbo"

# ---- FUNCTIONS ----
def extract_text_from_pdf(file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file.read())
        tmp_path = tmp.name
    doc = fitz.open(tmp_path)
    text = "\n".join([page.get_text() for page in doc])
    doc.close()
    os.remove(tmp_path)
    return text

def extract_text_from_docx(file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(file.read())
        tmp_path = tmp.name
    doc = docx.Document(tmp_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    os.remove(tmp_path)
    return text

def extract_text_from_pptx(file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(file.read())
        tmp_path = tmp.name
    prs = Presentation(tmp_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    os.remove(tmp_path)
    return text

def extract_text_from_excel(file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
        tmp.write(file.read())
        tmp_path = tmp.name
    df = pd.read_excel(tmp_path)
    text = df.to_string(index=False)
    os.remove(tmp_path)
    return text

def chat_with_gpt(messages):
    response = openai.chat.completions.create(
        model=MODEL,
        messages=messages,
        temperature=0.4,
        max_tokens=1000
    )
    return response.choices[0].message.content.strip()

# ---- EXPORT FUNCTIONS ----
def export_to_pdf(content, filename="ai_response.pdf"):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for line in content.split("\n"):
        pdf.multi_cell(0, 10, line)
    pdf.output(filename)
    return filename

def export_to_docx(content, filename="ai_response.docx"):
    doc = Document()
    doc.add_paragraph(content)
    doc.save(filename)
    return filename

def export_to_pptx(content, filename="ai_response.pptx"):
    prs = PptPresentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title, body = slide.shapes.title, slide.placeholders[1]
    title.text = "AI Response"
    body.text = content[:1000]  # Limited content
    prs.save(filename)
    return filename

def export_to_excel(content, filename="ai_response.xlsx"):
    df = pd.DataFrame({"Response": [content]})
    df.to_excel(filename, index=False)
    return filename

# ---- AUTHENTICATION ----
def check_password():
    def password_entered():
        if st.session_state["password"] == st.secrets["password"]:
            st.session_state["authenticated"] = True
        else:
            st.session_state["authenticated"] = False
            st.error("❌ Incorrect password")

    if "authenticated" not in st.session_state:
        st.text_input("Password", type="password", on_change=password_entered, key="password")
        st.stop()
    elif not st.session_state["authenticated"]:
        st.text_input("Password", type="password", on_change=password_entered, key="password")
        st.stop()

check_password()

# ---- UI ----
st.title("Lecturer AI Assistant")

uploaded_files = st.file_uploader("Upload PDFs, Word (DOCX), PPTX or Excel files", type=["pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)
user_query = st.text_input("Type your message")

if "doc_text" not in st.session_state:
    st.session_state.doc_text = ""
if "chat_history" not in st.session_state:
    st.session_state.chat_history = [
        {"role": "system", "content": "You are an assistant for a lecturer. Answer clearly and helpfully."}
    ]

if uploaded_files:
    all_text = []
    with st.spinner("Extracting text from uploaded files..."):
        for uploaded_file in uploaded_files:
            if uploaded_file.name.endswith(".pdf"):
                all_text.append(extract_text_from_pdf(uploaded_file))
            elif uploaded_file.name.endswith(".docx"):
                all_text.append(extract_text_from_docx(uploaded_file))
            elif uploaded_file.name.endswith(".pptx"):
                all_text.append(extract_text_from_pptx(uploaded_file))
            elif uploaded_file.name.endswith(".xlsx"):
                all_text.append(extract_text_from_excel(uploaded_file))
    st.session_state.doc_text = "\n\n".join(all_text)
    st.success(f"Loaded {len(uploaded_files)} file(s) successfully!")

if user_query:
    st.session_state.chat_history.append({"role": "user", "content": user_query})
    context_message = {"role": "user", "content": f"Context:\n{st.session_state.doc_text}"} if st.session_state.doc_text else None
    messages = st.session_state.chat_history.copy()
    if context_message:
        messages.insert(1, context_message)
    with st.spinner("Thinking..."):
        response = chat_with_gpt(messages)
        st.session_state.chat_history.append({"role": "assistant", "content": response})

# ---- Chat Display ----
if len(st.session_state.chat_history) > 1:
    st.markdown("### 💬 Conversation")
    for msg in st.session_state.chat_history[1:]:
        if msg["role"] == "user":
            st.markdown(f"**👤 You:** {msg['content']}")
        elif msg["role"] == "assistant":
            st.markdown(f"**🤖 AI:** {msg['content']}")

    st.markdown("---")
    st.markdown("### Export Last AI Response")

    last_response = st.session_state.chat_history[-1]["content"]
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        if st.button("Export to PDF"):
            pdf_file = export_to_pdf(last_response)
            with open(pdf_file, "rb") as f:
                st.download_button("Download PDF", f, file_name=pdf_file, mime="application/pdf")
    with col2:
        if st.button("Export to DOCX"):
            docx_file = export_to_docx(last_response)
            with open(docx_file, "rb") as f:
                st.download_button("Download DOCX", f, file_name=docx_file, mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    with col3:
        if st.button("Export to PPTX"):
            pptx_file = export_to_pptx(last_response)
            with open(pptx_file, "rb") as f:
                st.download_button("Download PPTX", f, file_name=pptx_file, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    with col4:
        if st.button("Export to Excel"):
            excel_file = export_to_excel(last_response)
            with open(excel_file, "rb") as f:
                st.download_button("Download Excel", f, file_name=excel_file, mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
